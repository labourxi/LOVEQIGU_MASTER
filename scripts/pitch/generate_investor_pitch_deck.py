#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""Generate AR游伴 investor pitch deck (.pptx) with eastern-space master theme."""

from __future__ import annotations

import os
from dataclasses import dataclass
from typing import Callable, List, Optional, Tuple

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = os.path.abspath(os.path.join(os.path.dirname(__file__), "../.."))
OUTPUT = os.path.join(ROOT, "docs/product/pitch/AR_YOUBAN_INVESTOR_PITCH_DECK_V1.pptx")

# Brand palette
INK = RGBColor(0x26, 0x3A, 0x34)
PAPER = RGBColor(0xF4, 0xF1, 0xEB)
GOLD = RGBColor(0xB6, 0x8A, 0x3D)
GRAY = RGBColor(0x6B, 0x72, 0x80)
LINE = RGBColor(0xC9, 0xB8, 0x96)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_INK = RGBColor(0x3D, 0x52, 0x4B)

FONT_TITLE = "Microsoft YaHei"
FONT_BODY = "Microsoft YaHei"
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


@dataclass
class SlideSpec:
    title: str
    body_lines: List[str]
    conclusion: str
    layout: str = "content"  # content | cover | closing | dark
    footnote: str = ""
    speaker_notes: str = ""


def rgb_hex(color: RGBColor) -> str:
    return f"#{color}"


def add_rect(slide, left, top, width, height, fill, line_color=None, line_width=0):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_width or 1)
    else:
        shape.line.fill.background()
    return shape


def add_round_rect(slide, left, top, width, height, fill, line_color=None):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_oval(slide, left, top, size, fill, line_color=None):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_line(slide, x1, y1, x2, y2, color=LINE, width=1.0):
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    conn.line.color.rgb = color
    conn.line.width = Pt(width)
    return conn


def set_text(
    shape,
    text: str,
    size: int = 18,
    bold: bool = False,
    color: RGBColor = INK,
    align=PP_ALIGN.LEFT,
    font=FONT_BODY,
):
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.name = font
    p.font.color.rgb = color
    p.alignment = align
    tf.word_wrap = True
    return tf


def add_textbox(slide, left, top, width, height, text, **kwargs):
    box = slide.shapes.add_textbox(left, top, width, height)
    set_text(box, text, **kwargs)
    return box


def add_multiline_box(slide, left, top, width, height, lines, size=16, color=INK, spacing=1.15):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.name = FONT_BODY
        p.font.color.rgb = color
        p.space_after = Pt(6)
        p.line_spacing = spacing
    tf.word_wrap = True
    return box


def draw_star_chain(slide, left, top, count=3, dark=False):
    dot_fill = GOLD if not dark else LINE
    line_color = LINE if not dark else RGBColor(0x8A, 0x9A, 0x92)
    gap = Inches(0.35)
    dot = Inches(0.08)
    for i in range(count):
        x = left + i * gap
        add_oval(slide, x, top, dot, dot_fill)
        if i < count - 1:
            add_line(slide, x + dot, top + dot / 2, x + gap, top + dot / 2, line_color, 0.75)


def draw_master_frame(slide, dark=False):
    bg = INK if dark else PAPER
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, bg)
    # top accent
    accent = GOLD if not dark else LINE
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), accent)
    # bottom meridian arc hint
    add_line(slide, Inches(0.8), Inches(6.95), Inches(12.5), Inches(6.95), accent if dark else LINE, 0.75)
    draw_star_chain(slide, Inches(11.2), Inches(0.35), 3, dark=dark)
    # brand mark
    mark_color = LINE if dark else GRAY
    add_textbox(
        slide,
        Inches(0.65),
        Inches(6.72),
        Inches(3),
        Inches(0.35),
        "AR游伴 · ARYOUBAN",
        size=9,
        color=mark_color,
    )


def add_conclusion_footer(slide, text: str, dark=False):
    color = LINE if dark else GRAY
    add_textbox(
        slide,
        Inches(0.75),
        Inches(6.35),
        Inches(11.8),
        Inches(0.45),
        f"结论 · {text}",
        size=11,
        color=color,
        align=PP_ALIGN.RIGHT,
    )


def add_speaker_notes(slide, notes: str):
    if not notes:
        return
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = notes


def add_title_block(slide, title: str, dark=False):
    color = PAPER if dark else INK
    add_textbox(slide, Inches(0.75), Inches(0.55), Inches(11.5), Inches(0.9), title, size=28, bold=True, color=color)
    add_line(slide, Inches(0.75), Inches(1.42), Inches(3.2), Inches(1.42), GOLD, 1.5)


def slide_cover(prs: Presentation, spec: SlideSpec):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    draw_master_frame(slide, dark=True)
    add_oval(slide, Inches(5.9), Inches(1.6), Inches(1.5), RGBColor(0x35, 0x4A, 0x44), LINE)
    add_textbox(slide, Inches(1.2), Inches(2.35), Inches(11), Inches(1.2), "AR游伴", size=54, bold=True, color=PAPER, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(1.2), Inches(3.35), Inches(11), Inches(0.8), "看见即是找回", size=32, color=GOLD, align=PP_ALIGN.CENTER)
    add_line(slide, Inches(4.2), Inches(4.25), Inches(9.1), Inches(4.25), LINE, 1)
    add_textbox(
        slide,
        Inches(1.5),
        Inches(4.55),
        Inches(10.3),
        Inches(0.6),
        "空间信物记忆系统 · 文旅数字化基础设施",
        size=18,
        color=LINE,
        align=PP_ALIGN.CENTER,
    )
    add_textbox(slide, Inches(1.2), Inches(5.35), Inches(11), Inches(0.5), "LOVEQIGU / ARYOUBAN · 2026", size=14, color=GRAY, align=PP_ALIGN.CENTER)
    add_conclusion_footer(slide, spec.conclusion, dark=True)
    add_speaker_notes(slide, spec.speaker_notes)


def slide_closing(prs: Presentation, spec: SlideSpec):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    draw_master_frame(slide, dark=True)
    add_textbox(slide, Inches(1.2), Inches(2.0), Inches(11), Inches(0.9), "看见即是找回", size=40, bold=True, color=PAPER, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(1.5), Inches(2.95), Inches(10.3), Inches(0.55), "让真实世界留下可被记住的探索痕迹", size=20, color=LINE, align=PP_ALIGN.CENTER)
    lines = [
        "联系人：[姓名]",
        "邮箱：[email@company.com]",
        "微信：[扫码添加]",
    ]
    add_multiline_box(slide, Inches(4.5), Inches(4.0), Inches(4.3), Inches(1.4), lines, size=16, color=PAPER)
    draw_star_chain(slide, Inches(5.8), Inches(5.55), 3, dark=True)
    add_conclusion_footer(slide, spec.conclusion, dark=True)
    add_speaker_notes(slide, spec.speaker_notes)


def slide_content(prs: Presentation, spec: SlideSpec):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    draw_master_frame(slide)
    add_title_block(slide, spec.title)
    add_multiline_box(slide, Inches(0.85), Inches(1.75), Inches(11.6), Inches(4.35), spec.body_lines, size=17, color=INK)
    if spec.footnote:
        add_textbox(slide, Inches(0.85), Inches(5.95), Inches(11), Inches(0.35), spec.footnote, size=10, color=GRAY)
    add_conclusion_footer(slide, spec.conclusion)
    add_speaker_notes(slide, spec.speaker_notes)


def slide_three_cards(prs: Presentation, spec: SlideSpec, cards: List[Tuple[str, List[str]]]):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    draw_master_frame(slide)
    add_title_block(slide, spec.title)
    card_w = Inches(3.55)
    gap = Inches(0.35)
    start_x = Inches(0.85)
    top = Inches(2.0)
    card_h = Inches(2.6)
    for i, (title, bullets) in enumerate(cards):
        x = start_x + i * (card_w + gap)
        add_round_rect(slide, x, top, card_w, card_h, WHITE, LINE)
        add_textbox(slide, x + Inches(0.2), top + Inches(0.2), card_w - Inches(0.4), Inches(0.45), title, size=16, bold=True, color=INK)
        add_multiline_box(slide, x + Inches(0.2), top + Inches(0.7), card_w - Inches(0.4), card_h - Inches(0.9), bullets, size=14, color=GRAY)
    if spec.body_lines:
        add_multiline_box(slide, Inches(0.85), Inches(4.85), Inches(11.5), Inches(1.2), spec.body_lines, size=15, color=INK)
    add_conclusion_footer(slide, spec.conclusion)
    add_speaker_notes(slide, spec.speaker_notes)


def slide_flow_steps(prs: Presentation, spec: SlideSpec, steps: List[Tuple[str, str, str]]):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    draw_master_frame(slide)
    add_title_block(slide, spec.title)
    n = len(steps)
    start_x = Inches(0.7)
    step_w = Inches(11.9) / n
    y = Inches(2.35)
    for i, (num, label, sub) in enumerate(steps):
        cx = start_x + step_w * i + step_w / 2
        add_oval(slide, cx - Inches(0.14), y, Inches(0.28), GOLD if i == 0 else LINE, GOLD if i == 0 else None)
        add_textbox(slide, cx - Inches(0.5), y + Inches(0.35), Inches(1.0), Inches(0.3), num, size=11, color=GRAY, align=PP_ALIGN.CENTER)
        add_textbox(slide, start_x + step_w * i + Inches(0.05), y + Inches(0.65), step_w - Inches(0.1), Inches(0.55), label, size=14, bold=True, color=INK, align=PP_ALIGN.CENTER)
        add_textbox(slide, start_x + step_w * i + Inches(0.05), y + Inches(1.2), step_w - Inches(0.1), Inches(0.55), sub, size=12, color=GRAY, align=PP_ALIGN.CENTER)
        if i < n - 1:
            add_line(slide, cx + Inches(0.2), y + Inches(0.14), cx + step_w - Inches(0.2), y + Inches(0.14), LINE, 1)
    if spec.body_lines:
        add_round_rect(slide, Inches(0.85), Inches(4.2), Inches(11.6), Inches(1.5), WHITE, LINE)
        add_multiline_box(slide, Inches(1.05), Inches(4.4), Inches(11.2), Inches(1.2), spec.body_lines, size=15, color=INK)
    add_conclusion_footer(slide, spec.conclusion)
    add_speaker_notes(slide, spec.speaker_notes)


def slide_table(prs: Presentation, spec: SlideSpec, headers: List[str], rows: List[List[str]], highlight_col: Optional[int] = None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    draw_master_frame(slide)
    add_title_block(slide, spec.title)
    cols = len(headers)
    left = Inches(0.75)
    top = Inches(1.85)
    width = Inches(11.8)
    row_h = Inches(0.42)
    col_w = width / cols
    # header
    add_rect(slide, left, top, width, row_h, INK)
    for c, h in enumerate(headers):
        add_textbox(slide, left + col_w * c + Inches(0.08), top + Inches(0.06), col_w, row_h, h, size=13, bold=True, color=PAPER)
    for r, row in enumerate(rows):
        y = top + row_h * (r + 1)
        bg = PAPER if r % 2 == 0 else WHITE
        add_rect(slide, left, y, width, row_h, bg, LINE)
        for c, cell in enumerate(row):
            color = GOLD if highlight_col is not None and c == highlight_col else INK
            bold = highlight_col is not None and c == highlight_col
            add_textbox(slide, left + col_w * c + Inches(0.08), y + Inches(0.05), col_w, row_h, cell, size=12, color=color, bold=bold)
    add_conclusion_footer(slide, spec.conclusion)
    add_speaker_notes(slide, spec.speaker_notes)


def slide_pilot_map(prs: Presentation, spec: SlideSpec):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    draw_master_frame(slide)
    add_title_block(slide, spec.title)
    cx, cy = Inches(6.65), Inches(2.55)
    add_round_rect(slide, cx - Inches(1.6), cy - Inches(0.35), Inches(3.2), Inches(0.7), INK)
    add_textbox(slide, cx - Inches(1.5), cy - Inches(0.28), Inches(3.0), Inches(0.55), "爱企谷（示例景区）", size=16, bold=True, color=PAPER, align=PP_ALIGN.CENTER)
    points = [
        ("点位 A", "入口仪式", Inches(2.2), Inches(4.2)),
        ("点位 B", "核心信物", cx - Inches(0.9), Inches(4.85)),
        ("点位 C", "商家礼遇", Inches(9.8), Inches(4.2)),
    ]
    for name, role, x, y in points:
        add_line(slide, cx, cy + Inches(0.35), x + Inches(0.9), y, LINE, 1)
        add_round_rect(slide, x, y, Inches(1.8), Inches(1.0), WHITE, GOLD)
        add_textbox(slide, x + Inches(0.12), y + Inches(0.12), Inches(1.56), Inches(0.35), name, size=14, bold=True, color=INK, align=PP_ALIGN.CENTER)
        add_textbox(slide, x + Inches(0.12), y + Inches(0.48), Inches(1.56), Inches(0.35), role, size=12, color=GRAY, align=PP_ALIGN.CENTER)
    metrics = "验证指标：完成率 · 信物获得率 · 权益领取率 · 复访率 · 传播素材产出"
    add_textbox(slide, Inches(0.85), Inches(5.85), Inches(11.5), Inches(0.4), metrics, size=14, color=INK, align=PP_ALIGN.CENTER)
    add_conclusion_footer(slide, spec.conclusion)
    add_speaker_notes(slide, spec.speaker_notes)


def slide_funding(prs: Presentation, spec: SlideSpec):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    draw_master_frame(slide)
    add_title_block(slide, spec.title)
    add_round_rect(slide, Inches(0.85), Inches(1.85), Inches(5.2), Inches(1.6), WHITE, GOLD)
    add_multiline_box(
        slide,
        Inches(1.05),
        Inches(2.05),
        Inches(4.8),
        Inches(1.3),
        ["轮次：种子轮 / Pre-A", "金额：[    ] 万元", "估值：[    ] 万元"],
        size=16,
        color=INK,
    )
    bars = [
        ("40%", "产品与技术", "XR Runtime · 后台 · 内容工厂"),
        ("30%", "景区试点与市场", "1→3 景区复制"),
        ("20%", "内容与运营", "信物生产 · 活动交付"),
        ("10%", "合规与基础设施", ""),
    ]
    top = Inches(1.95)
    for i, (pct, label, sub) in enumerate(bars):
        y = top + Inches(0.55) * i
        add_rect(slide, Inches(6.4), y, Inches(0.55), Inches(0.38), GOLD)
        add_textbox(slide, Inches(6.45), y + Inches(0.04), Inches(0.5), Inches(0.3), pct, size=11, bold=True, color=PAPER, align=PP_ALIGN.CENTER)
        add_textbox(slide, Inches(7.05), y, Inches(2.2), Inches(0.35), label, size=14, bold=True, color=INK)
        if sub:
            add_textbox(slide, Inches(9.2), y + Inches(0.02), Inches(3.2), Inches(0.35), sub, size=12, color=GRAY)
    add_conclusion_footer(slide, spec.conclusion)
    add_speaker_notes(slide, spec.speaker_notes)


def build_slides() -> List[Callable[[Presentation], None]]:
    builders: List[Callable[[Presentation], None]] = []

    builders.append(lambda prs: slide_cover(prs, SlideSpec(
        title="AR游伴",
        body_lines=[],
        conclusion="AR游伴不是 AR 展示工具，而是让真实空间留下可被记住的探索痕迹。",
        layout="cover",
        speaker_notes="开场 30 秒：先定义品类——空间信物记忆系统，不是 AR 特效公司。",
    )))

    builders.append(lambda prs: slide_content(prs, SlideSpec(
        title="文旅体验的三重断裂",
        body_lines=[
            "游客：打卡一次性，离开即失联",
            "景区：导览静态页，现场无记忆层",
            "商家：核销断链，难以承接探索行为",
            "",
            "         游客 ──╳── 景区 ──╳── 商家",
            "                    │",
            "                    ▼",
            "           数据孤岛 · 无记忆 · 无关系 · 无复访",
        ],
        conclusion="游客走完一圈，景区和商家却拿不到可持续的关系资产。",
        speaker_notes="用景区方痛点开场：游客来了，关系没留下。强调「关系资产」不是 DAU。",
    )))

    builders.append(lambda prs: slide_content(prs, SlideSpec(
        title="数字化做了「上线」，没做成「在场」",
        body_lines=[
            "传统路径：小程序 → 图文介绍 → 优惠券 → 结束",
            "",
            "缺失层（AR游伴补全）：",
            "  · 空间触发",
            "  · 仪式引导",
            "  · 信物记忆",
        ],
        conclusion="文旅数字化的缺口不在内容多少，而在空间现场是否产生可被记录的关系。",
    )))

    builders.append(lambda prs: slide_content(prs, SlideSpec(
        title="万亿文旅市场，正在从「导览」走向「在场操作系统」",
        body_lines=[
            "TAM · 中国文旅消费（万亿级市场）",
            "SAM · 景区 / 文博 / 城市文化点位数字化",
            "SOM · 空间信物 + XR 运营基础设施",
            "",
            "                    ▼",
            "              AR游伴切入点",
        ],
        footnote="市场规模分层为路演框架，具体数值需以最新行业公开报告核验。",
        conclusion="我们切入的不是 AR 特效市场，而是文旅场域数字化中「缺失的空间记忆层」。",
    )))

    builders.append(lambda prs: slide_three_cards(prs, SlideSpec(
        title="三个确定性趋势叠加",
        body_lines=["景区需要可运营、可复购、可传播的现场体验"],
        conclusion="政策、消费与超级 App 生态同时为「空间运营系统」打开了窗口期。",
    ), [
        ("政策推动", ["文化数字化", "文博活化"]),
        ("线下复苏", ["体验经济", "深度游回归"]),
        ("微信生态", ["小程序渗透", "超级入口"]),
    ]))

    builders.append(lambda prs: slide_content(prs, SlideSpec(
        title="AR游伴 = 现实空间的信物记忆操作系统",
        body_lines=[
            "┌─────────────────────────────────────┐",
            "│   空间操作系统层                     │",
            "│   触发 · 仪式 · 信物 · 关系          │",
            "└──────────────┬──────────────────────┘",
            "       XR 感知层  │  事件驱动层  │  运营回流层",
            "       识别/显现  │  Event Bus  │  活动/权益/CRM",
        ],
        conclusion="我们把现实空间变成可记录、可运营、可复访的体验结构，而非一次性打卡页面。",
        speaker_notes="核心 Slide：用「空间操作系统」定位，对标 OS 而非 App。",
    )))

    builders.append(lambda prs: slide_three_cards(prs, SlideSpec(
        title="XR 体验 × 信物系统 × 运营 CRM",
        body_lines=["用户端小程序统一承载三层能力", "信物层：故事进度资产 · 非数字藏品"],
        conclusion="产品不是单点 AR，而是现场体验、记忆资产与商业运营的一体化结构。",
    ), [
        ("XR 层", ["空间识别", "显现仪式", "备用显现"]),
        ("信物层", ["故事进度资产", "探索图谱节点", "回响记录"]),
        ("运营层", ["景区后台", "商家门户", "权益中心"]),
    ]))

    builders.append(lambda prs: slide_flow_steps(prs, SlideSpec(
        title="信物：用户与空间发生关系的记录",
        body_lines=[
            "信物 ≠ 数字藏品",
            "信物 = 故事进度资产  ·  藏品 = 传播/营销资产",
        ],
        conclusion="信物不是奖励道具，而是可被记住、可被运营、可被回看的空间关系凭证。",
        speaker_notes="必讲：信物≠数字藏品。信物是故事进度资产，藏品是传播资产。",
    ), [
        ("①", "空间触发", "XR入口"),
        ("②", "仪式引导", "探索状态"),
        ("③", "信物显现", "收藏册"),
        ("④", "意义沉淀", "今日回响"),
    ]))

    builders.append(lambda prs: slide_flow_steps(prs, SlideSpec(
        title="进入 → 探索 → 发现 → 获得信物 → 完成",
        body_lines=[],
        conclusion="用户走完一轮，获得的是一段可被记住的探索体验，而不是一个积分数字。",
    ), [
        ("①", "进入景区", "xr_start_v1"),
        ("②", "探索地图", "space_trail_v1"),
        ("③", "发现印记", "AR显现"),
        ("④", "显现信物", "relic_emerge_v1"),
        ("⑤", "完成回响", "权益/复访"),
    ]))

    builders.append(lambda prs: slide_content(prs, SlideSpec(
        title="游客感知：克制、东方、在场",
        body_lines=[
            "视觉：星图 · 经络 · 宣纸 · 淡金",
            "交互：探索地图 · 权益中心 · 个人印记",
            "语言：探索 / 显现 / 回响 / 礼遇",
            "边界：无等级 · 无抽卡 · 无战力",
        ],
        conclusion="我们刻意远离游戏化，以东方空间美学建立文旅场景的信任感与高级感。",
    )))

    builders.append(lambda prs: slide_content(prs, SlideSpec(
        title="事件驱动的 XR 运行时，UI 与 XR 解耦",
        body_lines=[
            "用户 UI 层：首页 · 探索地图 · AR入口 · 信物册 · 权益中心",
            "           ↓ safeNavigate / safeTap",
            "Event Bus：XR_USER_TRIGGER · RELIC_CREATED · STAR_LIGHTED",
            "     ↓                              ↓",
            "XR Runtime（识别/管线/降级）   World Engine（结构化世界建模）",
        ],
        conclusion="技术壁垒不在单个 AR 特效，而在可复用、可降级、可运营的空间运行时。",
    )))

    builders.append(lambda prs: slide_content(prs, SlideSpec(
        title="真机 XR 优先，备用显现保障交付",
        body_lines=[
            "用户触发 → AR 显现 ──成功──→ 信物流程",
            "              │",
            "              │ 设备/权限/环境不支持",
            "              ▼",
            "         备用显现 ──完成──→ 同一信物闭环",
        ],
        conclusion="景区现场网络与设备不可控，系统必须保证体验闭环不断裂。",
    )))

    builders.append(lambda prs: slide_content(prs, SlideSpec(
        title="B2B2C：景区付费，商家增值，用户免费探索",
        body_lines=[
            "① 景区 SaaS 年费/项目费     → 景区/文博     空间系统部署",
            "② 点位内容生产费             → 景区+平台     信物/显现内容",
            "③ 商家活动服务费             → 在地商家     权益投放/核销",
            "④ 增值运营模块               → 景区         数据分析/活动引擎",
            "⑤ 城市复制授权               → 政府/文旅集团 多景区复制",
            "",
            "用户端：免费进入探索 · 信物与权益自然获得",
        ],
        conclusion="收入来自空间运营基础设施，而非向游客售卖虚拟道具。",
        speaker_notes="商业模式：B2B2C，景区买单。用户免费探索，避免 C 端道具争议。",
    )))

    builders.append(lambda prs: slide_pilot_map(prs, SlideSpec(
        title="最小可验证单元：一个景区，三个探索点",
        body_lines=[],
        conclusion="一个景区、三个点位，即可跑通从进入到信物到权益的完整商业闭环。",
        speaker_notes="试点可验证：90 天、3 点位、5 项核心指标。可现场演示小程序闭环。",
    )))

    builders.append(lambda prs: slide_three_cards(prs, SlideSpec(
        title="90 天试点节奏（示意）",
        body_lines=[],
        conclusion="试点目标不是 Demo，而是可复制的景区交付模板与运营数据基线。",
    ), [
        ("Day 0–30", ["点位勘测", "XR 标定", "信物定义"]),
        ("Day 31–60", ["内容上线", "试运行", "商家权益接入"]),
        ("Day 61–90", ["运营复盘", "复制模板冻结", "二景区 BD 启动"]),
    ]))

    builders.append(lambda prs: slide_table(prs, SlideSpec(
        title="每一次探索，都是可运营的关系数据",
        body_lines=[],
        conclusion="信物是用户资产，行为链是景区 CRM——两者在同一系统内闭环。",
    ), ["行为事件", "结构化输出", "运营动作"], [
        ["进入景区 / 探索点", "探索路径图谱", "路径优化"],
        ["AR显现 / 备用完成", "显现成功率", "设备策略"],
        ["信物获得", "信物档案", "内容迭代"],
        ["权益领取", "礼遇转化", "商家协同"],
        ["复访 / 分享", "传播触点", "活动投放"],
    ]))

    builders.append(lambda prs: slide_three_cards(prs, SlideSpec(
        title="平台 · 景区 · 商家，同一套数据语法",
        body_lines=[],
        conclusion="数据不是报表堆砌，而是驱动景区复购与商家转化的运营语言。",
    ), [
        ("平台总览", ["审查/发布", "内容生产线"]),
        ("景区看板", ["探索进度", "信物/活动"]),
        ("商家门户", ["核销/礼遇", "活动效果"]),
    ]))

    builders.append(lambda prs: slide_content(prs, SlideSpec(
        title="四维壁垒",
        body_lines=[
            "              空间操作系统定位",
            "                     │",
            "   世界观一致性 ──┼── 工程可交付",
            "   (Canon冻结)    │    (XR降级)",
            "                     │",
            "   商业闭环完整 ──┴── 非游戏化品牌信任",
            "   (信物→权益)",
        ],
        conclusion="竞品多做 AR 特效或打卡工具，我们做的是景区可长期运营的空间记忆基础设施。",
    )))

    builders.append(lambda prs: slide_table(prs, SlideSpec(
        title="AR游伴 vs 传统方案",
        body_lines=[],
        conclusion="我们不是更炫的 AR，而是更完整的景区空间运营系统。",
    ), ["维度", "传统导览小程序", "AR特效外包", "AR游伴"], [
        ["空间触发", "弱", "中", "强"],
        ["信物记忆", "无", "无", "核心"],
        ["运营后台", "弱", "无", "完整"],
        ["商业闭环", "券-only", "项目制", "SaaS+运营"],
        ["可复制性", "低", "低", "高"],
        ["品牌气质", "工具感", "娱乐感", "文旅高级感"],
    ], highlight_col=3))

    builders.append(lambda prs: slide_flow_steps(prs, SlideSpec(
        title="从单景区模板，到城市文化网络",
        body_lines=["○ ─── ○ ─── ○ ─── ○ ─── ○   城市星图网络"],
        conclusion="复制的是空间操作系统与内容工厂，不是一次性定制项目。",
    ), [
        ("P1", "1景区·3点位", "试点验证"),
        ("P2", "1城市·N景区", "标准化复制"),
        ("P3", "城际信物图谱", "文化IP网络"),
    ]))

    builders.append(lambda prs: slide_content(prs, SlideSpec(
        title="传播飞轮 × 商业飞轮",
        body_lines=[
            "传播飞轮：",
            "  震撼显现 → 用户拍摄分享 → 新客进入 → 新信物 → 再传播",
            "",
            "商业飞轮：",
            "  探索完成 → 信物获得 → 权益领取 → 商家转化 → 景区复购",
            "",
            "           ╭────────────────╮",
            "           │  AR游伴核心     │",
            "           │  空间信物系统   │",
            "           ╰────────────────╯",
        ],
        conclusion="传播带来新客，商业带来复购，信物系统是两者的交汇点。",
    )))

    builders.append(lambda prs: slide_table(prs, SlideSpec(
        title="四方共赢",
        body_lines=[],
        conclusion="游客获得记忆，景区获得关系，商家获得客流，平台获得规模。",
    ), ["角色", "价值 1", "价值 2", "价值 3"], [
        ["游客", "被记住的探索体验", "文化共鸣", "礼遇延续"],
        ["景区", "可运营现场", "数据资产", "差异化 IP"],
        ["商家", "精准到店", "权益转化", "低门槛参与"],
        ["平台", "标准化复制", "内容审查发布", "城市级扩展"],
    ]))

    builders.append(lambda prs: slide_funding(prs, SlideSpec(
        title="融资轮次与用途（示意框架）",
        body_lines=[],
        conclusion="融资用于把「已验证的单景区闭环」推进到「可复制的城市级空间网络」。",
        speaker_notes="融资用途按 40/30/20/10 结构讲清。金额与估值由创始团队现场填入。",
    )))

    builders.append(lambda prs: slide_flow_steps(prs, SlideSpec(
        title="从试点到规模 · 18 个月路线图",
        body_lines=["每一步都可交付、可验收、可复盘，不以概念换估值。"],
        conclusion="每一步都可交付、可验收、可复盘，不以概念换估值。",
    ), [
        ("Q1", "单景区试点", "信物闭环冻结"),
        ("Q2", "3景区复制", "商家权益规模化"),
        ("Q3", "城市签约", "内容工厂上线"),
        ("Q4", "标准化SaaS", "数据产品化"),
    ]))

    builders.append(lambda prs: slide_closing(prs, SlideSpec(
        title="看见即是找回",
        body_lines=[],
        conclusion="我们在寻找认同「空间记忆」长期价值的文旅与科技同行者。",
        layout="closing",
        speaker_notes="收尾邀请：欢迎现场体验「进入景区」完整闭环。留下联系方式。",
    )))

    return builders


def generate():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    for builder in build_slides():
        builder(prs)
    os.makedirs(os.path.dirname(OUTPUT), exist_ok=True)
    prs.save(OUTPUT)
    return OUTPUT


if __name__ == "__main__":
    path = generate()
    print(f"PITCH_DECK_GENERATED: {path}")
    print(f"SLIDES: 25")
